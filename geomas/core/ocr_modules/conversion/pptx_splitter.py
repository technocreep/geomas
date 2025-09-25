from __future__ import annotations

import hashlib
import logging
from io import BytesIO
from pathlib import Path
from typing import Callable, List

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..io.fs import atomic_write
from ..io.paths import content_addressed_work_path
from .base import Converter
from .text_table_to_pdf import TextTableToPDF


class PPTXSplitter(Converter):
    """Extract text and images from PowerPoint files."""

    def __init__(self, work_dir: Path, *, redact: bool = False) -> None:
        self.work_dir = work_dir
        self.redact = redact
        self.text_converter = TextTableToPDF(work_dir)

    def _export_image(self, blob: bytes) -> Path:
        digest = hashlib.sha256(blob).hexdigest()
        out = self.work_dir / "images" / f"{digest}.png"
        with Image.open(BytesIO(blob)) as img:
            exif = None if self.redact else img.info.get("exif")
            buf = BytesIO()
            img.save(buf, format="PNG", exif=exif)
        atomic_write(out, buf.getvalue())
        return out

    def convert(
        self,
        path: Path,
        *,
        get_logger: Callable[[str, str | None], logging.LoggerAdapter],
        request_id: str | None = None,
    ) -> tuple[Path, List[Path]]:
        prs = Presentation(str(path))
        text_tmp = content_addressed_work_path(self.work_dir, path, ".txt")
        lines = [
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        image_paths = [
            self._export_image(shape.image.blob)
            for slide in prs.slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        text_tmp.write_text("\n".join(lines), encoding="utf-8")
        pdf_path = self.text_converter.convert(
            text_tmp, get_logger=get_logger, request_id=request_id
        )
        text_tmp.unlink(missing_ok=True)
        return pdf_path, image_paths


__all__ = ["PPTXSplitter"]
